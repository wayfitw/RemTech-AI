"""Генерация документов.
- create_docx: markdown-подобный текст → .docx (портировано из mybot _create_docx_sync)
- create_pdf:  markdown-подобный текст → .pdf (reportlab, кириллица через TTF из конфига)
"""
import datetime as dt
import html
import io
import os
import re
import sys

from app.config import get_settings

# TTF с кириллицей для PDF: из конфига, иначе системный дефолт по ОС.
PDF_FONT_PATH = get_settings().pdf_font_path or (
    "C:/Windows/Fonts/arial.ttf" if sys.platform == "win32"
    else "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf"
)


def create_docx(content: str, filename: str = "document", letterhead: bool = True) -> bytes:
    """letterhead=True — фирменный бланк «Ремтехники» вверху (RT + реквизиты);
    False — классический стиль без бланка (по выбору пользователя для договоров/док-тов)."""
    from docx import Document
    from docx.enum.table import WD_ALIGN_VERTICAL
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.oxml import OxmlElement
    from docx.oxml.ns import qn
    from docx.shared import Cm, Pt, RGBColor

    from services.docx_style import COMPANY

    GRAPHITE = RGBColor(0x26, 0x28, 0x2F)

    doc = Document()
    doc.core_properties.title = filename

    section = doc.sections[0]
    section.page_width = Cm(21)
    section.page_height = Cm(29.7)
    section.left_margin = Cm(2.0)
    section.right_margin = Cm(1.5)
    section.top_margin = Cm(0.46)
    section.bottom_margin = Cm(1.0)

    normal = doc.styles["Normal"]
    normal.font.name = "Times New Roman"
    normal.font.size = Pt(12)
    normal.paragraph_format.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY
    normal.paragraph_format.space_after = Pt(0)
    normal.paragraph_format.line_spacing = 1.15

    RE_SECTION = re.compile(r"^(\d{1,2})\.\s+[А-ЯЁA-Z]")
    RE_SUBSECTION = re.compile(r"^\d+\.\d+")
    RE_TITLE_CAPS = re.compile(r"^[А-ЯЁA-Z\s№«»\d\-–—.,]+$")
    RE_MD_TABLE_ROW = re.compile(r"^\|(.+)\|$")
    RE_MD_TABLE_SEP = re.compile(r"^\|[-| :]+\|$")

    def _set_font(run, bold=False, size_pt=12, color=None):
        run.font.name = "Times New Roman"
        run.font.size = Pt(size_pt)
        run.bold = bold
        if color:
            run.font.color.rgb = color

    def _add_run(para, text, bold=False, size_pt=12, color=None):
        run = para.add_run(text)
        _set_font(run, bold=bold, size_pt=size_pt, color=color)
        return run

    def _add_inline(para, text):
        parts = text.split("**")
        for i, part in enumerate(parts):
            if part:
                run = para.add_run(part)
                _set_font(run, bold=(i % 2 == 1))

    def _set_cell_border(cell, color="D9D9D9"):
        tc = cell._tc
        tcPr = tc.get_or_add_tcPr()
        borders = OxmlElement("w:tcBorders")
        for side in ("top", "left", "bottom", "right"):
            tag = OxmlElement(f"w:{side}")
            tag.set(qn("w:val"), "single")
            tag.set(qn("w:sz"), "4")
            tag.set(qn("w:color"), color)
            borders.append(tag)
        tcPr.append(borders)

    def _set_cell_no_border(cell):
        tc = cell._tc
        tcPr = tc.get_or_add_tcPr()
        tcBorders = OxmlElement("w:tcBorders")
        for side in ("top", "left", "bottom", "right", "insideH", "insideV"):
            tag = OxmlElement(f"w:{side}")
            tag.set(qn("w:val"), "none")
            tcBorders.append(tag)
        tcPr.append(tcBorders)

    def _set_cell_shading(cell, fill="D9D9D9"):
        tc = cell._tc
        tcPr = tc.get_or_add_tcPr()
        shd = OxmlElement("w:shd")
        shd.set(qn("w:val"), "clear")
        shd.set(qn("w:color"), "auto")
        shd.set(qn("w:fill"), fill)
        tcPr.append(shd)

    def _para_in_cell(cell, text, bold=False, align=WD_ALIGN_PARAGRAPH.LEFT, size_pt=12, color=None):
        p = cell.paragraphs[0] if cell.paragraphs else cell.add_paragraph()
        p.clear()
        p.alignment = align
        p.paragraph_format.space_after = Pt(0)
        _add_run(p, text, bold=bold, size_pt=size_pt, color=color)
        return p

    def _yellow_rule(p):
        """Тонкая фирменная жёлтая линия под абзацем (акцент под главным заголовком)."""
        pPr = p._p.get_or_add_pPr()
        pbdr = OxmlElement("w:pBdr")
        bottom = OxmlElement("w:bottom")
        bottom.set(qn("w:val"), "single")
        bottom.set(qn("w:sz"), "12")
        bottom.set(qn("w:space"), "4")
        bottom.set(qn("w:color"), "FFCB05")
        pbdr.append(bottom)
        pPr.append(pbdr)

    def _build_header(logo_path, req_lines):
        hdr = doc.sections[0].header
        hdr.is_linked_to_previous = False
        for p in hdr.paragraphs:
            p._element.getparent().remove(p._element)
        htbl = hdr.add_table(rows=1, cols=2, width=Cm(17.5))
        htbl.autofit = False
        left_cell, right_cell = htbl.rows[0].cells[0], htbl.rows[0].cells[1]
        left_cell.width = Cm(6.0)
        right_cell.width = Cm(11.5)
        for cell in (left_cell, right_cell):
            _set_cell_no_border(cell)
        lp = left_cell.paragraphs[0]
        lp.paragraph_format.space_after = Pt(0)
        if logo_path and os.path.exists(logo_path):
            lp.add_run().add_picture(logo_path, width=Cm(5.7))
        else:
            _set_font(lp.add_run("[ЛОГОТИП]"), size_pt=11)
        for idx, line in enumerate(req_lines):
            rp = right_cell.paragraphs[0] if idx == 0 else right_cell.add_paragraph()
            rp.alignment = WD_ALIGN_PARAGRAPH.RIGHT
            rp.paragraph_format.space_after = Pt(0)
            _set_font(rp.add_run(line), size_pt=11)

    def _make_md_table(rows):
        # Фирменный стиль: тёмная шапка + белый текст, чередование строк, тонкие линии.
        WHITE = RGBColor(0xFF, 0xFF, 0xFF)
        ncols = max(len(r) for r in rows)
        tbl = doc.add_table(rows=len(rows), cols=ncols)
        for i, row in enumerate(rows):
            is_header = i == 0
            for j in range(ncols):
                cell = tbl.rows[i].cells[j]
                _para_in_cell(
                    cell, (row[j].strip() if j < len(row) else ""), bold=is_header,
                    align=WD_ALIGN_PARAGRAPH.CENTER if is_header else WD_ALIGN_PARAGRAPH.LEFT,
                    color=WHITE if is_header else None,
                )
                if is_header:
                    _set_cell_shading(cell, "2B2E33")        # графит-шапка
                elif i % 2 == 0:
                    _set_cell_shading(cell, "F4F4F6")        # приглушённое чередование
                _set_cell_border(cell)                        # тонкая линия D9D9D9
        return tbl

    def _make_2col_table(rows, widths=(50, 50), borders=False):
        tbl = doc.add_table(rows=len(rows), cols=2)
        for i, (left, right) in enumerate(rows):
            lc, rc = tbl.rows[i].cells[0], tbl.rows[i].cells[1]
            _para_in_cell(lc, left)
            _para_in_cell(rc, right)
            for cell in (lc, rc):
                (_set_cell_border if borders else _set_cell_no_border)(cell)
        total_cm = 16.0
        for row in tbl.rows:
            row.cells[0].width = Cm(total_cm * widths[0] / 100)
            row.cells[1].width = Cm(total_cm * widths[1] / 100)
        return tbl

    def _letterhead():
        """Фирменный бланк вверху: жёлтый бокс «RT» + ООО «Ремтехника» + ИНН/КПП/ОГРН.
        Логотип рисуется ячейкой (файл не нужен). На первой странице документа."""
        tbl = doc.add_table(rows=1, cols=2)
        tbl.autofit = False
        lc, rc = tbl.rows[0].cells
        lc.width, rc.width = Cm(1.7), Cm(15.3)
        tbl.rows[0].height = Cm(1.3)
        for cell in (lc, rc):
            _set_cell_no_border(cell)
            cell.vertical_alignment = WD_ALIGN_VERTICAL.CENTER
        # левая ячейка — жёлтый бокс с «RT»
        _set_cell_shading(lc, "FFCB05")
        lp = lc.paragraphs[0]
        lp.alignment = WD_ALIGN_PARAGRAPH.CENTER
        lp.paragraph_format.space_after = Pt(0)
        _set_font(lp.add_run("RT"), bold=True, size_pt=22, color=RGBColor(0x1A, 0x1A, 0x1A))
        # правая ячейка — название + реквизиты
        grey = RGBColor(0x7F, 0x7F, 0x7F)
        rows_txt = [
            (COMPANY["name"], True, GRAPHITE, 13),
            (f"ИНН {COMPANY['inn']}   КПП {COMPANY['kpp']}", False, grey, 10),
            (f"ОГРН {COMPANY['ogrn']}", False, grey, 10),
        ]
        for idx, (t, bold, color, sz) in enumerate(rows_txt):
            rp = rc.paragraphs[0] if idx == 0 else rc.add_paragraph()
            rp.alignment = WD_ALIGN_PARAGRAPH.LEFT
            rp.paragraph_format.space_after = Pt(0)
            rp.paragraph_format.line_spacing = 1.0
            _set_font(rp.add_run(t), bold=bold, size_pt=sz, color=color)
        # небольшой отступ под бланком (пустой абзац — редактор доков его пропускает)
        doc.add_paragraph().paragraph_format.space_after = Pt(6)

    # Фирменный бланк — если включён и пользователь не задал свой [HEADER].
    if letterhead and "[HEADER" not in content.upper():
        _letterhead()

    # Декодируем HTML-сущности (&nbsp;, &amp;, …), которые модель иногда вставляет,
    # иначе они выводятся буквально; неразрывный пробел → обычный.
    content = html.unescape(content).replace("\xa0", " ")
    lines = content.split("\n")
    i = 0
    while i < len(lines):
        line = lines[i].strip()

        if line.upper().startswith("[HEADER"):
            m = re.search(r"logo=([^\]]*)", line, re.I)
            logo_path = m.group(1).strip() if m else ""
            req_lines = []
            i += 1
            while i < len(lines) and not lines[i].strip().upper().startswith("[/HEADER"):
                req_lines.append(lines[i].strip())
                i += 1
            _build_header(logo_path, req_lines)
            i += 1
            continue

        if line.upper() == "[PAGEBREAK]":
            doc.add_page_break()
            i += 1
            continue

        if line.upper().startswith("[2COL"):
            m = re.search(r"widths?=(\d+)[,/](\d+)", line, re.I)
            w = (int(m.group(1)), int(m.group(2))) if m else (50, 50)
            rows_2col = []
            i += 1
            while i < len(lines) and not lines[i].strip().upper().startswith("[/2COL"):
                parts = lines[i].split("||")
                rows_2col.append((parts[0].strip() if parts else "",
                                  parts[1].strip() if len(parts) > 1 else ""))
                i += 1
            _make_2col_table(rows_2col, widths=w)
            i += 1
            continue

        if RE_MD_TABLE_ROW.match(line):
            md_rows = []
            while i < len(lines):
                l = lines[i].strip()
                if RE_MD_TABLE_SEP.match(l):
                    i += 1
                    continue
                if not RE_MD_TABLE_ROW.match(l):
                    break
                md_rows.append([c.strip() for c in l.strip("|").split("|")])
                i += 1
            if md_rows:
                _make_md_table(md_rows)
            continue

        i += 1

        if not line or line in ("---", "—--", "***"):
            doc.add_paragraph("")
            continue

        if line.startswith("# "):
            p = doc.add_paragraph(); p.alignment = WD_ALIGN_PARAGRAPH.CENTER
            _add_run(p, line[2:], bold=True, size_pt=14, color=GRAPHITE)
            _yellow_rule(p); continue
        if line.startswith("## "):
            p = doc.add_paragraph(); p.alignment = WD_ALIGN_PARAGRAPH.LEFT
            p.paragraph_format.space_before = Pt(10)
            _add_run(p, line[3:], bold=True, size_pt=13, color=GRAPHITE); continue
        if line.startswith("### "):
            p = doc.add_paragraph(); p.alignment = WD_ALIGN_PARAGRAPH.LEFT
            p.paragraph_format.space_before = Pt(8)
            _add_run(p, line[4:], bold=True, size_pt=12, color=GRAPHITE); continue

        if RE_SECTION.match(line) and not RE_SUBSECTION.match(line):
            p = doc.add_paragraph(); p.alignment = WD_ALIGN_PARAGRAPH.LEFT
            p.paragraph_format.space_before = Pt(10)
            _add_run(p, line, bold=True, size_pt=12, color=GRAPHITE); continue

        if RE_TITLE_CAPS.match(line) and 2 <= len(line) <= 80 and len(line.split()) <= 12:
            p = doc.add_paragraph(); p.alignment = WD_ALIGN_PARAGRAPH.CENTER
            _add_run(p, line, bold=True, size_pt=13, color=GRAPHITE); continue

        p = doc.add_paragraph()
        if RE_SUBSECTION.match(line) or line.startswith("—") or line.startswith("-"):
            p.paragraph_format.left_indent = Cm(0.63)
        else:
            p.paragraph_format.first_line_indent = Cm(0.63)
        _add_inline(p, line)

    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


def _num(x) -> str:
    """Число с разделителями тысяч: 1234567 → '1 234 567'."""
    try:
        return f"{round(float(x)):,}".replace(",", " ")
    except Exception:
        return str(x)


def create_proposal(data: dict) -> bytes:
    """Коммерческое предложение (КП) в фирменном стиле РемТехники (Word).

    data: {title, client, executor, contact, validity_days, notes, markup_percent,
    items: [{name, qty, price}]}. price — базовая цена за единицу; сумма считается
    с наценкой markup_percent.
    """
    from docx import Document
    from docx.enum.table import WD_TABLE_ALIGNMENT
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.oxml import OxmlElement
    from docx.oxml.ns import qn
    from docx.shared import Pt, RGBColor

    from services.docx_style import (  # общий стиль (issue #19)
        BAND,
        HAIRLINE,
        INK,
        SOFT,
        YELLOW,
        requisites_lines,
        shade,
    )

    ink = RGBColor(0x2B, 0x2E, 0x33)   # графит — заголовки/шапка/итог
    grey = RGBColor(0x80, 0x80, 0x80)  # серый — второстепенный текст

    doc = Document()
    normal = doc.styles["Normal"]
    normal.font.name = "Calibri"
    normal.font.size = Pt(11)

    def _bottom_rule(p, color_hex, sz=22):
        """Тонкая линия-акцент под абзацем (нижняя граница)."""
        pbdr = OxmlElement("w:pBdr")
        b = OxmlElement("w:bottom")
        b.set(qn("w:val"), "single"); b.set(qn("w:sz"), str(sz))
        b.set(qn("w:space"), "6"); b.set(qn("w:color"), color_hex)
        pbdr.append(b)
        p._p.get_or_add_pPr().append(pbdr)

    def _thin_borders(tbl, color_hex=HAIRLINE):
        """Мягкая тонкая сетка вместо чёрного Table Grid."""
        borders = OxmlElement("w:tblBorders")
        for edge in ("top", "left", "bottom", "right", "insideH", "insideV"):
            e = OxmlElement(f"w:{edge}")
            e.set(qn("w:val"), "single"); e.set(qn("w:sz"), "4")
            e.set(qn("w:space"), "0"); e.set(qn("w:color"), color_hex)
            borders.append(e)
        tbl._tbl.tblPr.append(borders)

    # ── Заголовок с тонким жёлтым акцентом (вместо крупной плашки) ────────────
    tp = doc.add_paragraph()
    tp.paragraph_format.space_after = Pt(2)
    r = tp.add_run("Коммерческое предложение")
    r.bold = True
    r.font.size = Pt(22)
    r.font.color.rgb = ink
    _bottom_rule(tp, YELLOW)   # фирменный акцент — тонкая жёлтая линия под заголовком

    executor = data.get("executor") or "ООО «Ремтехника»"
    meta = doc.add_paragraph()
    meta.paragraph_format.space_before = Pt(4)
    mr = meta.add_run(f"{executor}   ·   {dt.datetime.now():%d.%m.%Y}")
    mr.font.size = Pt(10)
    mr.font.color.rgb = grey

    if data.get("client"):
        p = doc.add_paragraph()
        p.paragraph_format.space_after = Pt(0)
        p.add_run("Кому: ").bold = True
        p.add_run(str(data["client"]))
    if data.get("title"):
        p = doc.add_paragraph()
        tr = p.add_run(str(data["title"]))
        tr.bold = True
        tr.font.size = Pt(13)
        tr.font.color.rgb = ink

    # ── Таблица позиций ──────────────────────────────────────────────────────
    items = data.get("items") or []
    markup = float(data.get("markup_percent") or 0)
    headers = ["№", "Наименование", "Кол-во", "Цена, ₽", "Сумма, ₽"]
    aligns = [WD_ALIGN_PARAGRAPH.CENTER, WD_ALIGN_PARAGRAPH.LEFT, WD_ALIGN_PARAGRAPH.CENTER,
              WD_ALIGN_PARAGRAPH.RIGHT, WD_ALIGN_PARAGRAPH.RIGHT]
    t = doc.add_table(rows=1, cols=len(headers))
    t.alignment = WD_TABLE_ALIGNMENT.CENTER
    _thin_borders(t)
    for i, h in enumerate(headers):
        pp = t.rows[0].cells[i].paragraphs[0]
        pp.alignment = aligns[i]
        run = pp.add_run(h)
        run.bold = True
        run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        run.font.size = Pt(10)
        shade(t.rows[0].cells[i], INK)

    total = 0
    for idx, it in enumerate(items, 1):
        qty = float(it.get("qty") or 1)
        price = float(it.get("price") or 0)
        summ = round(qty * price * (1 + markup / 100))
        total += summ
        vals = [str(idx), str(it.get("name", "")), _num(qty), _num(price), _num(summ)]
        cells = t.add_row().cells
        for i, v in enumerate(vals):
            pp = cells[i].paragraphs[0]
            pp.alignment = aligns[i]
            run = pp.add_run(str(v))
            run.font.size = Pt(10)
            if idx % 2 == 0:
                shade(cells[i], SOFT)   # приглушённое серое чередование

    # ИТОГО — сдержанно: серая подпись + бледно-жёлтая сумма (единственный акцент)
    trow = t.add_row().cells
    tl = trow[0].merge(trow[3])
    tl.paragraphs[0].alignment = WD_ALIGN_PARAGRAPH.RIGHT
    lab = tl.paragraphs[0].add_run("ИТОГО" + (f" (с наценкой {markup:g}%)" if markup else ""))
    lab.bold = True
    lab.font.color.rgb = ink
    trow[4].paragraphs[0].alignment = WD_ALIGN_PARAGRAPH.RIGHT
    tot = trow[4].paragraphs[0].add_run(_num(total) + " ₽")
    tot.bold = True
    tot.font.color.rgb = ink
    shade(tl, SOFT)
    shade(trow[4], BAND)   # бледно-жёлтый — мягкий фирменный акцент

    # ── Условия и контакты ───────────────────────────────────────────────────
    doc.add_paragraph()
    if data.get("validity_days"):
        doc.add_paragraph(f"Предложение действительно {int(data['validity_days'])} рабочих дней.")
    if data.get("notes"):
        doc.add_paragraph(str(data["notes"]))
    if data.get("contact"):
        p = doc.add_paragraph()
        p.add_run("Контакт: ").bold = True
        p.add_run(str(data["contact"]))

    # ── Реквизиты компании (issue #26) ───────────────────────────────────────
    sep = doc.add_paragraph()
    sep.paragraph_format.space_before = Pt(10)
    _bottom_rule(sep, HAIRLINE, sz=6)   # тонкий разделитель над реквизитами
    for line in requisites_lines():
        rp = doc.add_paragraph()
        rp.paragraph_format.space_after = Pt(0)
        rr = rp.add_run(line)
        rr.font.size = Pt(9)
        rr.font.color.rgb = grey

    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


def create_spec_report(data: dict) -> bytes:
    """Issue #25 — отчёт анализа технического задания (Word, фирменный стиль).

    data: {title, summary, requirements:[...], risks:[...], contradictions:[...],
    gaps:[...]}. Разбор ТЗ делает модель; здесь — оформление результата."""
    from docx import Document
    from docx.shared import Pt, RGBColor

    from services.docx_style import DARK, YELLOW, shade

    ink = RGBColor(0x1A, 0x1A, 0x1A)
    grey = RGBColor(0x7F, 0x7F, 0x7F)
    doc = Document()
    doc.styles["Normal"].font.name = "Calibri"
    doc.styles["Normal"].font.size = Pt(11)

    tbar = doc.add_table(rows=1, cols=1)
    bc = tbar.rows[0].cells[0]
    shade(bc, YELLOW)
    r = bc.paragraphs[0].add_run("АНАЛИЗ ТЕХНИЧЕСКОГО ЗАДАНИЯ")
    r.bold = True
    r.font.size = Pt(16)
    r.font.color.rgb = ink

    if data.get("title"):
        p = doc.add_paragraph()
        tr = p.add_run(str(data["title"]))
        tr.bold = True
        tr.font.size = Pt(13)
    m = doc.add_paragraph().add_run(f"Сформировано {dt.datetime.now():%d.%m.%Y}")
    m.font.size = Pt(9)
    m.font.color.rgb = grey

    if data.get("summary"):
        doc.add_paragraph(str(data["summary"]))

    # Сводная таблица «категория → количество»
    sections = [
        ("Требования", data.get("requirements") or []),
        ("Риски", data.get("risks") or []),
        ("Противоречия", data.get("contradictions") or []),
        ("Пробелы (не хватает)", data.get("gaps") or []),
    ]
    t = doc.add_table(rows=1, cols=2)
    t.style = "Table Grid"
    for i, h in enumerate(("Категория", "Количество")):
        cell = t.rows[0].cells[i]
        rr = cell.paragraphs[0].add_run(h)
        rr.bold = True
        rr.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        shade(cell, DARK)
    for label, items in sections:
        cells = t.add_row().cells
        cells[0].text = label
        cells[1].text = str(len(items))

    for label, items in sections:
        if not items:
            continue
        h = doc.add_paragraph()
        hr = h.add_run(label)
        hr.bold = True
        hr.font.size = Pt(13)
        for it in items:
            doc.add_paragraph(str(it), style="List Bullet")

    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


def create_conversation_report(data: dict) -> bytes:
    """Issue #41 (TASK-0502) — отчёт анализа звонка/переписки (Word, фирменный стиль).

    data: {title, summary, agreements:[...], open_questions:[...], next_steps:[...],
    risks:[...]}. Разбор материала делает модель; здесь — оформление результата."""
    from docx import Document
    from docx.shared import Pt, RGBColor

    from services.docx_style import DARK, YELLOW, shade

    ink = RGBColor(0x1A, 0x1A, 0x1A)
    grey = RGBColor(0x7F, 0x7F, 0x7F)
    doc = Document()
    doc.styles["Normal"].font.name = "Calibri"
    doc.styles["Normal"].font.size = Pt(11)

    tbar = doc.add_table(rows=1, cols=1)
    bc = tbar.rows[0].cells[0]
    shade(bc, YELLOW)
    r = bc.paragraphs[0].add_run("АНАЛИЗ ПЕРЕПИСКИ / ЗВОНКА")
    r.bold = True
    r.font.size = Pt(16)
    r.font.color.rgb = ink

    if data.get("title"):
        p = doc.add_paragraph()
        tr = p.add_run(str(data["title"]))
        tr.bold = True
        tr.font.size = Pt(13)
    m = doc.add_paragraph().add_run(f"Сформировано {dt.datetime.now():%d.%m.%Y}")
    m.font.size = Pt(9)
    m.font.color.rgb = grey

    if data.get("summary"):
        doc.add_paragraph(str(data["summary"]))

    sections = [
        ("Договорённости", data.get("agreements") or []),
        ("Открытые вопросы", data.get("open_questions") or []),
        ("Следующие шаги (кому / что / когда)", data.get("next_steps") or []),
        ("Риски", data.get("risks") or []),
    ]
    t = doc.add_table(rows=1, cols=2)
    t.style = "Table Grid"
    for i, h in enumerate(("Раздел", "Пунктов")):
        cell = t.rows[0].cells[i]
        rr = cell.paragraphs[0].add_run(h)
        rr.bold = True
        rr.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        shade(cell, DARK)
    for label, items in sections:
        cells = t.add_row().cells
        cells[0].text = label
        cells[1].text = str(len(items))

    for label, items in sections:
        if not items:
            continue
        h = doc.add_paragraph()
        hr = h.add_run(label)
        hr.bold = True
        hr.font.size = Pt(13)
        for it in items:
            doc.add_paragraph(str(it), style="List Bullet")

    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


def create_estimate(data: dict) -> bytes:
    """Issue #27 — смета/бюджет в Excel (.xlsx) с формулами и фирменным оформлением.

    data: {title, client, notes, markup_percent, items:[{name, unit, qty, price}]}.
    Суммы и итог — настоящие Excel-формулы (лист остаётся редактируемым)."""
    from openpyxl import Workbook
    from openpyxl.styles import Alignment, Border, Font, PatternFill, Side

    from services.docx_style import BAND, DARK, YELLOW

    wb = Workbook()
    ws = wb.active
    ws.title = "Смета"

    thin = Side(style="thin", color="BFBFBF")
    box = Border(left=thin, right=thin, top=thin, bottom=thin)
    center = Alignment(horizontal="center", vertical="center")
    left = Alignment(horizontal="left", vertical="center", wrap_text=True)
    money = "# ##0"

    headers = ["№", "Наименование", "Ед.", "Кол-во", "Цена, ₽", "Сумма, ₽"]
    ncols = len(headers)

    # Титул
    ws.merge_cells(start_row=1, start_column=1, end_row=1, end_column=ncols)
    c = ws.cell(1, 1, (data.get("title") or "Смета").upper())
    c.font = Font(bold=True, size=14, color=DARK)
    c.fill = PatternFill("solid", fgColor=YELLOW)
    c.alignment = Alignment(horizontal="left", vertical="center")
    ws.row_dimensions[1].height = 26
    row = 2
    if data.get("client"):
        ws.cell(row, 1, f"Заказчик: {data['client']}").font = Font(size=10)
        row += 1
    row += 1

    # Шапка таблицы
    head_row = row
    for i, h in enumerate(headers, 1):
        cell = ws.cell(head_row, i, h)
        cell.font = Font(bold=True, color="FFFFFF")
        cell.fill = PatternFill("solid", fgColor=DARK)
        cell.alignment = center
        cell.border = box
    row += 1

    markup = float(data.get("markup_percent") or 0)
    factor = 1 + markup / 100
    first = row
    for idx, it in enumerate(data.get("items") or [], 1):
        qty = float(it.get("qty") or 1)
        price = float(it.get("price") or 0)
        vals = [idx, str(it.get("name", "")), str(it.get("unit") or "шт"), qty, price]
        for i, v in enumerate(vals, 1):
            cell = ws.cell(row, i, v)
            cell.border = box
            cell.alignment = left if i == 2 else center
            if i in (4, 5):
                cell.number_format = money
        # Сумма — формула: Кол-во * Цена * (1 + наценка)
        s = ws.cell(row, 6, f"=ROUND(D{row}*E{row}*{factor},0)")
        s.border = box
        s.number_format = money
        s.alignment = center
        if idx % 2 == 0:
            for i in range(1, ncols + 1):
                ws.cell(row, i).fill = PatternFill("solid", fgColor=BAND)
        row += 1
    last = row - 1

    # ИТОГО — формула SUM
    ws.merge_cells(start_row=row, start_column=1, end_row=row, end_column=5)
    tl = ws.cell(row, 1, "ИТОГО" + (f" (с наценкой {markup:g}%)" if markup else ""))
    tl.font = Font(bold=True)
    tl.alignment = Alignment(horizontal="right", vertical="center")
    tl.fill = PatternFill("solid", fgColor=YELLOW)
    total = ws.cell(row, 6, f"=SUM(F{first}:F{last})" if last >= first else 0)
    total.font = Font(bold=True)
    total.number_format = money
    total.fill = PatternFill("solid", fgColor=YELLOW)
    total.border = box
    tl.border = box
    row += 2

    if data.get("notes"):
        ws.cell(row, 1, str(data["notes"])).font = Font(size=10, italic=True)

    for col, w in zip("ABCDEF", (5, 42, 8, 10, 14, 16)):
        ws.column_dimensions[col].width = w

    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


_PLACEHOLDER_RE = re.compile(r"\{\{([^}]+)\}\}")


def fill_template(data: bytes, values: dict) -> tuple[bytes, list[str], list[str]]:
    """Issue #26 — заполнение фирменного шаблона .docx: подстановка {{ПОЛЕ}} → значение
    с сохранением форматирования. Возвращает (bytes, заполненные_поля, оставшиеся_плейсхолдеры)."""
    from docx import Document

    doc = Document(io.BytesIO(data))
    filled: set[str] = set()

    def _replace(text: str) -> str:
        def sub(m):
            key = m.group(1).strip()
            if key in values:
                filled.add(key)
                return str(values[key])
            return m.group(0)   # оставляем незаполненным
        return _PLACEHOLDER_RE.sub(sub, text)

    def _process(par) -> None:
        full = "".join(r.text for r in par.runs)
        if "{{" not in full:
            return
        new = _replace(full)
        if new != full and par.runs:
            # переписываем в первый run (сохраняет стиль абзаца), остальные очищаем
            par.runs[0].text = new
            for r in par.runs[1:]:
                r.text = ""

    for par in doc.paragraphs:
        _process(par)
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                for par in cell.paragraphs:
                    _process(par)

    buf = io.BytesIO()
    doc.save(buf)
    # оставшиеся незаполненные плейсхолдеры (по всему тексту)
    all_text = "\n".join(p.text for p in Document(io.BytesIO(buf.getvalue())).paragraphs)
    for t in Document(io.BytesIO(buf.getvalue())).tables:
        for row in t.rows:
            for cell in row.cells:
                all_text += "\n" + cell.text
    remaining = sorted(set(_PLACEHOLDER_RE.findall(all_text)))
    return buf.getvalue(), sorted(filled), remaining


def create_proposal_pdf(data: dict) -> bytes:
    """Коммерческое предложение (КП) в фирменном стиле «Ремтехники» (PDF, issue #28).
    Те же данные, что и create_proposal: {title, client, executor, contact,
    validity_days, notes, markup_percent, items:[{name, qty, price}]}."""
    from reportlab.lib import colors
    from reportlab.lib.enums import TA_LEFT
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.styles import ParagraphStyle
    from reportlab.lib.units import cm
    from reportlab.platypus import Paragraph, SimpleDocTemplate, Spacer, Table, TableStyle

    from services.docx_style import COMPANY, requisites_lines

    font = _register_pdf_font()
    yellow = colors.HexColor("#FFCB05")
    dark = colors.HexColor("#1A1A1A")
    band = colors.HexColor("#FFF6D5")
    grey = colors.HexColor("#7F7F7F")

    body = ParagraphStyle("B", fontName=font, fontSize=11, leading=15, alignment=TA_LEFT)
    small = ParagraphStyle("S", fontName=font, fontSize=8.5, leading=11, textColor=grey)
    title = ParagraphStyle("T", fontName=font, fontSize=17, leading=21, textColor=dark)

    buf = io.BytesIO()
    pdf = SimpleDocTemplate(buf, pagesize=A4, leftMargin=2 * cm, rightMargin=1.5 * cm,
                            topMargin=1.5 * cm, bottomMargin=1.5 * cm, title="КП")
    flow = []

    def esc(s):
        return str(s).replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")

    # Жёлтая титульная плашка
    bar = Table([[Paragraph("<b>КОММЕРЧЕСКОЕ ПРЕДЛОЖЕНИЕ</b>", title)]], colWidths=[17.5 * cm])
    bar.setStyle(TableStyle([("BACKGROUND", (0, 0), (-1, -1), yellow),
                             ("TOPPADDING", (0, 0), (-1, -1), 8), ("BOTTOMPADDING", (0, 0), (-1, -1), 8),
                             ("LEFTPADDING", (0, 0), (-1, -1), 10)]))
    flow.append(bar)
    executor = data.get("executor") or COMPANY["name"]
    flow.append(Paragraph(f"{esc(executor)} · {dt.datetime.now():%d.%m.%Y}", small))
    flow.append(Spacer(1, 8))
    if data.get("client"):
        flow.append(Paragraph(f"<b>Кому:</b> {esc(data['client'])}", body))
    if data.get("title"):
        flow.append(Paragraph(f"<b>{esc(data['title'])}</b>", body))
    flow.append(Spacer(1, 8))

    # Таблица позиций
    markup = float(data.get("markup_percent") or 0)
    rows = [["№", "Наименование", "Кол-во", "Цена, ₽", "Сумма, ₽"]]
    total = 0
    for idx, it in enumerate(data.get("items") or [], 1):
        qty = float(it.get("qty") or 1)
        price = float(it.get("price") or 0)
        summ = round(qty * price * (1 + markup / 100))
        total += summ
        rows.append([str(idx), Paragraph(esc(it.get("name", "")), body), _num(qty), _num(price), _num(summ)])
    total_label = "ИТОГО" + (f" (с наценкой {markup:g}%)" if markup else "")
    rows.append([total_label, "", "", "", _num(total) + " ₽"])

    tbl = Table(rows, colWidths=[1 * cm, 8.5 * cm, 2 * cm, 3 * cm, 3 * cm])
    style = [
        ("FONTNAME", (0, 0), (-1, -1), font), ("FONTSIZE", (0, 0), (-1, -1), 9),
        ("BACKGROUND", (0, 0), (-1, 0), dark), ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
        ("GRID", (0, 0), (-1, -1), 0.5, colors.HexColor("#BFBFBF")),
        ("VALIGN", (0, 0), (-1, -1), "MIDDLE"), ("TOPPADDING", (0, 0), (-1, -1), 4),
        ("BOTTOMPADDING", (0, 0), (-1, -1), 4),
        ("SPAN", (0, -1), (3, -1)), ("BACKGROUND", (0, -1), (-1, -1), yellow),
        ("FONTNAME", (0, -1), (-1, -1), font),
    ]
    for i in range(2, len(rows) - 1, 2):
        style.append(("BACKGROUND", (0, i), (-1, i), band))
    tbl.setStyle(TableStyle(style))
    flow.append(tbl)
    flow.append(Spacer(1, 10))

    if data.get("validity_days"):
        flow.append(Paragraph(f"Предложение действительно {int(data['validity_days'])} рабочих дней.", body))
    if data.get("notes"):
        flow.append(Paragraph(esc(data["notes"]), body))
    if data.get("contact"):
        flow.append(Paragraph(f"<b>Контакт:</b> {esc(data['contact'])}", body))
    flow.append(Spacer(1, 14))
    for line in requisites_lines():
        flow.append(Paragraph(esc(line), small))

    pdf.build(flow)
    return buf.getvalue()


_PDF_FONT_REGISTERED = False


def _register_pdf_font() -> str:
    """Регистрирует TTF с кириллицей в reportlab. Возвращает имя шрифта."""
    global _PDF_FONT_REGISTERED
    from reportlab.pdfbase import pdfmetrics
    from reportlab.pdfbase.ttfonts import TTFont

    if _PDF_FONT_REGISTERED:
        return "AppFont"
    if PDF_FONT_PATH and os.path.exists(PDF_FONT_PATH):
        pdfmetrics.registerFont(TTFont("AppFont", PDF_FONT_PATH))
        _PDF_FONT_REGISTERED = True
        return "AppFont"
    return "Helvetica"  # без кириллицы, но не падаем


def create_pdf(content: str, filename: str = "document") -> bytes:
    from reportlab.lib.enums import TA_CENTER, TA_JUSTIFY
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
    from reportlab.lib.units import cm
    from reportlab.platypus import Paragraph, SimpleDocTemplate, Spacer

    font = _register_pdf_font()
    styles = getSampleStyleSheet()
    body = ParagraphStyle("Body", parent=styles["Normal"], fontName=font,
                          fontSize=11, leading=15, alignment=TA_JUSTIFY)
    h1 = ParagraphStyle("H1", parent=body, fontSize=15, leading=19,
                        alignment=TA_CENTER, spaceAfter=10)
    h2 = ParagraphStyle("H2", parent=body, fontSize=13, leading=17,
                        alignment=TA_CENTER, spaceAfter=8)

    buf = io.BytesIO()
    pdf = SimpleDocTemplate(buf, pagesize=A4, leftMargin=2 * cm, rightMargin=1.5 * cm,
                            topMargin=1.5 * cm, bottomMargin=1.5 * cm, title=filename)
    flow = []

    def esc(s: str) -> str:
        s = s.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
        # **bold** → <b>
        s = re.sub(r"\*\*(.+?)\*\*", r"<b>\1</b>", s)
        return s

    for raw in content.split("\n"):
        line = raw.strip()
        if not line:
            flow.append(Spacer(1, 6))
        elif line.startswith("# "):
            flow.append(Paragraph(esc(line[2:]), h1))
        elif line.startswith("## "):
            flow.append(Paragraph(esc(line[3:]), h2))
        elif line.startswith("### "):
            flow.append(Paragraph(esc(line[4:]), h2))
        else:
            flow.append(Paragraph(esc(line), body))

    pdf.build(flow)
    return buf.getvalue()


def create_presentation(spec: dict) -> bytes:
    """Собирает современную презентацию (.pptx) в фирменном стиле «Ремтехники»
    (графит + жёлтый акцент, без логотипа). Поддерживает картинки на слайдах:
    обложка-фон, двухколоночные слайды (текст + фото), слайды-разделители.

    Картинки в spec приходят уже как БАЙТЫ (генерацию/поиск делает оркестратор —
    docgen сети не трогает):
        spec["_cover_image"]      — bytes фонового изображения обложки (опц.)
        slide["_image"]           — bytes картинки слайда (опц.)
        slide["layout"]           — "split" (по умолч., текст+картинка) | "full"
                                     (картинка на весь слайд + заголовок поверх)
    spec = {
        "title": "...", "subtitle": "...", "author": "Ремтехника",
        "slides": [{"title": "...", "bullets": [...], "notes": "...",
                    "_image": b"...", "layout": "split"}, ...],
    }
    Пустой список слайдов допустим (будет только титульный).
    """
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
    from pptx.oxml.ns import qn
    from pptx.util import Emu, Pt

    INK = RGBColor(0x2B, 0x2E, 0x33)     # графит — заголовки
    BODY = RGBColor(0x33, 0x36, 0x3B)    # основной текст
    GREY = RGBColor(0x7F, 0x7F, 0x7F)    # второстепенный
    YELLOW = RGBColor(0xFF, 0xCB, 0x05)  # фирменный акцент
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    DARK = RGBColor(0x1A, 0x1C, 0x20)    # скрим под белым текстом на фото
    FONT = "Calibri"

    EMU = 914400
    SW, SH = int(13.333 * EMU), int(7.5 * EMU)          # 16:9
    MX = int(0.9 * EMU)                                  # боковые поля

    prs = Presentation()
    prs.slide_width, prs.slide_height = SW, SH
    blank = prs.slide_layouts[6]

    def _bg(slide, color=WHITE):
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = color

    def _rect(slide, x, y, w, h, color, alpha=None):
        sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
        sp.fill.solid()
        sp.fill.fore_color.rgb = color
        sp.line.fill.background()
        sp.shadow.inherit = False
        if alpha is not None:   # alpha = доля НЕПРОЗРАЧНОСТИ, 0..1
            solid = sp._element.spPr.find(qn("a:solidFill"))
            srgb = solid.find(qn("a:srgbClr")) if solid is not None else None
            if srgb is not None:
                srgb.append(srgb.makeelement(qn("a:alpha"), {"val": str(int(alpha * 100000))}))
        return sp

    def _text(slide, x, y, w, h, anchor=MSO_ANCHOR.TOP):
        tb = slide.shapes.add_textbox(x, y, w, h)
        tf = tb.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = anchor
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        return tf

    def _runs(p, text, color=None, size=None):
        """Разбивает **жирный** на run'ы (как в остальных генераторах)."""
        for i, part in enumerate(str(text).split("**")):
            if not part:
                continue
            r = p.add_run()
            r.text, r.font.name = part, FONT
            r.font.bold = i % 2 == 1
            if color is not None:
                r.font.color.rgb = color
            if size is not None:
                r.font.size = Pt(size)

    def _cover_image(slide, data, x, y, w, h):
        """Вставляет картинку с обрезкой «под панель» (cover-fill, без искажения)."""
        try:
            from PIL import Image
            iw, ih = Image.open(io.BytesIO(data)).size
        except Exception:
            iw = ih = 0
        pic = slide.shapes.add_picture(io.BytesIO(data), x, y, width=w, height=h)
        if iw and ih:
            img_ar, panel_ar = iw / ih, w / h
            if img_ar > panel_ar:                       # шире панели — режем бока
                c = (1 - panel_ar / img_ar) / 2
                pic.crop_left = pic.crop_right = c
            else:                                        # выше панели — режем верх/низ
                c = (1 - img_ar / panel_ar) / 2
                pic.crop_top = pic.crop_bottom = c
        return pic

    def _footer(slide, author, idx, light=False):
        col = WHITE if light else GREY
        f = _text(slide, MX, SH - int(0.5 * EMU), SW - 2 * MX, int(0.3 * EMU))
        r = f.paragraphs[0].add_run()
        r.text, r.font.size, r.font.color.rgb, r.font.name = author, Pt(9), col, FONT
        n = _text(slide, SW - MX - int(0.6 * EMU), SH - int(0.5 * EMU), int(0.6 * EMU), int(0.3 * EMU))
        n.paragraphs[0].alignment = PP_ALIGN.RIGHT
        r = n.paragraphs[0].add_run()
        r.text, r.font.size, r.font.color.rgb, r.font.name = str(idx), Pt(9), col, FONT

    def _bullets(slide, bullets, x, y, w, h, color=BODY):
        btf = _text(slide, x, y, w, h)
        for i, b in enumerate(bullets):
            p = btf.paragraphs[0] if i == 0 else btf.add_paragraph()
            p.space_after = Pt(10)
            lead = p.add_run()
            lead.text, lead.font.name, lead.font.bold = "•  ", FONT, True
            lead.font.size, lead.font.color.rgb = Pt(18), YELLOW
            _runs(p, b, color=color, size=18)

    author = (spec.get("author") or "Ремтехника").strip()

    # ── титульный слайд ──────────────────────────────────────────────────────
    s0 = prs.slides.add_slide(blank)
    cover = spec.get("_cover_image")
    if cover:
        _bg(s0, DARK)
        _cover_image(s0, cover, 0, 0, SW, SH)             # фон на весь слайд
        _rect(s0, 0, 0, SW, SH, DARK, alpha=0.45)         # затемнение под текст
        _rect(s0, 0, 0, int(0.16 * EMU), SH, YELLOW)      # акцент слева
        tf = _text(s0, MX, int(2.5 * EMU), SW - 2 * MX, int(2.4 * EMU))
        _runs(tf.paragraphs[0], spec.get("title") or "Презентация", color=WHITE, size=44)
        for r in tf.paragraphs[0].runs:
            r.font.bold = True
        if spec.get("subtitle"):
            ps = tf.add_paragraph(); ps.space_before = Pt(14)
            _runs(ps, spec["subtitle"], color=RGBColor(0xE6, 0xE6, 0xE6), size=20)
        _footer(s0, author, 1, light=True)
    else:
        _bg(s0)
        _rect(s0, 0, 0, int(0.16 * EMU), SH, YELLOW)
        tf = _text(s0, MX, int(2.4 * EMU), SW - 2 * MX, int(2.4 * EMU))
        _runs(tf.paragraphs[0], spec.get("title") or "Презентация", color=INK, size=40)
        for r in tf.paragraphs[0].runs:
            r.font.bold = True
        if spec.get("subtitle"):
            ps = tf.add_paragraph(); ps.space_before = Pt(14)
            _runs(ps, spec["subtitle"], color=GREY, size=20)
        _footer(s0, author, 1)

    # ── контентные слайды ────────────────────────────────────────────────────
    for idx, sl in enumerate(spec.get("slides") or [], start=2):
        s = prs.slides.add_slide(blank)
        img = sl.get("_image")
        layout = (sl.get("layout") or "").lower()
        bullets = sl.get("bullets") or []
        title = sl.get("title") or ""

        if img and layout == "full":
            # слайд-разделитель: фото на весь слайд + заголовок поверх
            _bg(s, DARK)
            _cover_image(s, img, 0, 0, SW, SH)
            _rect(s, 0, 0, SW, SH, DARK, alpha=0.40)
            _rect(s, MX, int(3.05 * EMU), int(1.1 * EMU), Emu(34000), YELLOW)
            htf = _text(s, MX, int(3.3 * EMU), SW - 2 * MX, int(2.0 * EMU))
            _runs(htf.paragraphs[0], title, color=WHITE, size=34)
            for r in htf.paragraphs[0].runs:
                r.font.bold = True
            if bullets:
                _bullets(s, bullets[:3], MX, int(4.6 * EMU), SW - 2 * MX, int(2.0 * EMU),
                         color=RGBColor(0xEC, 0xEC, 0xEC))
            _footer(s, author, idx, light=True)
        elif img:
            # двухколоночный: текст слева, фото — правая колонка на всю высоту
            _bg(s)
            panel_w = int(SW * 0.42)
            panel_x = SW - panel_w
            _cover_image(s, img, panel_x, 0, panel_w, SH)
            _rect(s, panel_x - Emu(30000), 0, Emu(30000), SH, YELLOW)   # тонкая граница-акцент
            col_w = panel_x - MX - int(0.4 * EMU)
            htf = _text(s, MX, int(0.6 * EMU), col_w, int(1.0 * EMU))
            _runs(htf.paragraphs[0], title, color=INK, size=26)
            for r in htf.paragraphs[0].runs:
                r.font.bold = True
            _rect(s, MX, int(1.55 * EMU), int(1.1 * EMU), Emu(28000), YELLOW)
            if bullets:
                _bullets(s, bullets, MX, int(1.95 * EMU), col_w, int(4.7 * EMU))
            _footer(s, author, idx)
        else:
            # без картинки — заголовок + тезисы на всю ширину
            _bg(s)
            htf = _text(s, MX, int(0.6 * EMU), SW - 2 * MX, int(1.0 * EMU))
            _runs(htf.paragraphs[0], title, color=INK, size=28)
            for r in htf.paragraphs[0].runs:
                r.font.bold = True
            _rect(s, MX, int(1.55 * EMU), int(1.1 * EMU), Emu(28000), YELLOW)
            if bullets:
                _bullets(s, bullets, MX, int(1.95 * EMU), SW - 2 * MX, int(4.7 * EMU))
            _footer(s, author, idx)

        notes = (sl.get("notes") or "").strip()
        if notes:
            s.notes_slide.notes_text_frame.text = notes

    out = io.BytesIO()
    prs.save(out)
    return out.getvalue()


# TASK-0507 (#45) — КП-презентация на технику (PPTX). Вынесена в отдельный модуль
# (объёмная раскладка), доступна как docgen.create_proposal_pptx для симметрии с
# create_proposal / create_proposal_pdf.
from services.proposal_pptx import create_proposal_pptx  # noqa: E402,F401
