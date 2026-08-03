"""#45 Этап 3 — REST КП-презентаций: RBAC, безопасность фото (owner-only),
генерация + скачивание, извлечение (модель замокана)."""
import io

from PIL import Image

from services import docgen, proposal_pptx


def _auth(t):
    return {"Authorization": f"Bearer {t}"}


def _png():
    b = io.BytesIO()
    Image.new("RGB", (80, 60), (200, 160, 40)).save(b, "PNG")
    return b.getvalue()


async def _register_admin(client):
    r = await client.post("/api/register",
                          json={"username": "director", "password": "pass1234", "full_name": "Дир"})
    return r.json()["token"]


async def _make_user(client, admin, username, role):
    await client.post("/api/admin/users",
                      json={"username": username, "password": "pass1234", "full_name": username, "role": role},
                      headers=_auth(admin))
    return (await client.post("/api/login",
                              json={"username": username, "password": "pass1234"})).json()["token"]


async def test_proposals_rbac(client):
    admin = await _register_admin(client)
    worker = await _make_user(client, admin, "worker", "user")   # не продажи/руководство
    r = await client.post("/api/proposals/photo",
                          files={"file": ("p.png", _png(), "image/png")}, headers=_auth(worker))
    assert r.status_code == 403                                  # роль не допущена
    # полностью без авторизации (чистим cookie сессии, оставленный логином) → 401
    client.cookies.clear()
    assert (await client.post("/api/proposals/generate",
                              json={"name": "X", "blocks": []})).status_code == 401


async def test_proposals_photo_generate_download(client):
    admin = await _register_admin(client)                        # admin проходит RBAC
    up = await client.post("/api/proposals/photo",
                           files={"file": ("exc.png", _png(), "image/png")}, headers=_auth(admin))
    assert up.status_code == 200
    img_id = up.json()["image_id"]
    payload = {"name": "Экскаватор XCMG", "filename": "КП_XE215",
               "blocks": [{"type": "title", "title": "Экскаватор XCMG"},
                          {"type": "split", "image_id": img_id, "rows": [["Мощность", "118 кВт"]]}]}
    gen = await client.post("/api/proposals/generate", json=payload, headers=_auth(admin))
    assert gen.status_code == 200, gen.text
    fid = gen.json()["file_id"]
    dl = await client.get(f"/api/files/{fid}", headers=_auth(admin))
    assert dl.status_code == 200 and dl.content[:2] == b"PK"
    from pptx import Presentation
    prs = Presentation(io.BytesIO(dl.content))
    pics = sum(1 for s in prs.slides for sh in s.shapes if sh.shape_type == 13)
    assert pics == 1                                             # фото владельца встроено


async def test_proposals_photo_owner_only(client):
    # чужой image_id → 403 (белый список = только собственные файлы, критерий безопасности)
    admin = await _register_admin(client)
    sales = await _make_user(client, admin, "seller", "продажи")
    up = await client.post("/api/proposals/photo",
                           files={"file": ("a.png", _png(), "image/png")}, headers=_auth(admin))
    other_id = up.json()["image_id"]                            # файл принадлежит admin
    payload = {"name": "X", "blocks": [{"type": "split", "image_id": other_id, "rows": [["a", "b"]]}]}
    r = await client.post("/api/proposals/generate", json=payload, headers=_auth(sales))
    assert r.status_code == 403                                  # продажник не возьмёт чужое фото


async def test_proposals_extract(client, monkeypatch):
    admin = await _register_admin(client)

    async def fake_extract(data, filename, db, **kw):
        return {"name": "Погрузчик LiuGong CLG856H", "brand": "LiuGong",
                "blocks": [{"type": "title", "title": "Погрузчик LiuGong CLG856H"}]}

    monkeypatch.setattr(proposal_pptx, "extract_slides_from_document", fake_extract)
    docx = docgen.create_docx("Погрузчик LiuGong CLG856H. Цена 12 400 000 руб.", "kp")
    r = await client.post(
        "/api/proposals/extract",
        files={"file": ("kp.docx", docx,
                        "application/vnd.openxmlformats-officedocument.wordprocessingml.document")},
        headers=_auth(admin))
    assert r.status_code == 200, r.text
    assert r.json()["brand"] == "LiuGong" and r.json()["blocks"][0]["type"] == "title"


# ── #54 (TASK-0510): история последних 30 КП ─────────────────────────────────

async def _generate(client, admin, name="Экскаватор XCMG", client_name="ООО «Стройка»"):
    payload = {"name": name, "client_name": client_name, "filename": name,
               "blocks": [{"type": "title", "title": name}]}
    r = await client.post("/api/proposals/generate", json=payload, headers=_auth(admin))
    assert r.status_code == 200, r.text
    return r.json()


async def test_history_records_and_lists(client):
    admin = await _register_admin(client)
    gen = await _generate(client, admin)

    lst = await client.get("/api/proposals/history", headers=_auth(admin))
    assert lst.status_code == 200
    rows = lst.json()
    assert len(rows) == 1
    assert rows[0]["file_id"] == gen["file_id"]
    assert rows[0]["machine"] == "Экскаватор XCMG" and rows[0]["client_name"] == "ООО «Стройка»"

    # повторное скачивание по сохранённому file_id
    dl = await client.get(f"/api/files/{rows[0]['file_id']}", headers=_auth(admin))
    assert dl.status_code == 200 and dl.content[:2] == b"PK"

    # снимок контракта — основа для нового КП
    item = await client.get(f"/api/proposals/history/{rows[0]['id']}", headers=_auth(admin))
    assert item.status_code == 200
    assert item.json()["payload"]["name"] == "Экскаватор XCMG"


async def test_history_isolated_by_owner(client):
    admin = await _register_admin(client)
    gen = await _generate(client, admin)
    hist = (await client.get("/api/proposals/history", headers=_auth(admin))).json()
    item_id = hist[0]["id"]

    seller = await _make_user(client, admin, "seller3", "продажи")
    # чужая история не видна
    assert (await client.get("/api/proposals/history", headers=_auth(seller))).json() == []
    # чужая запись и чужой файл — 403
    assert (await client.get(f"/api/proposals/history/{item_id}",
                             headers=_auth(seller))).status_code == 403
    assert (await client.get(f"/api/files/{gen['file_id']}",
                             headers=_auth(seller))).status_code == 403


async def test_history_retention_keeps_last_30(session):
    """Ретенция: у пользователя остаются только 30 последних записей."""
    from app import repositories as repo
    u = await repo.create_user(session, "manager30", "h$1", role="продажи")
    await session.commit()
    for i in range(33):
        await repo.add_proposal_history(session, user_id=u.id, file_id=None,
                                        file_name=f"КП-{i}.pptx", machine=f"Модель {i}")
    await session.commit()
    rows = await repo.list_proposal_history(session, u.id, limit=100)
    assert len(rows) == repo.PROPOSAL_HISTORY_LIMIT == 30
    assert rows[0].file_name == "КП-32.pptx"          # новейшая первой
    assert all("КП-2.pptx" != r.file_name for r in rows)   # самые старые вычищены


# ── #55 (TASK-0511): отправка PPTX в Telegram ────────────────────────────────

async def test_send_telegram_success(client, monkeypatch):
    admin = await _register_admin(client)
    gen = await _generate(client, admin)

    import app.routers.proposals as pr
    from app.config import get_settings
    sent = {}

    class FakeTx:
        def __init__(self, token):
            pass

        async def send_file(self, chat_id, data, filename, method, field):
            sent.update(chat_id=chat_id, name=filename, size=len(data), method=method)

        async def aclose(self):
            pass

    settings = get_settings()
    monkeypatch.setattr(settings, "telegram_bot_token", "test-token", raising=False)
    monkeypatch.setattr(settings, "telegram_allowlist", "555:director", raising=False)
    monkeypatch.setattr(pr, "get_settings", lambda: settings)
    import app.telegram_bot as tb
    monkeypatch.setattr(tb, "TelegramTransport", FakeTx)

    r = await client.post("/api/proposals/send-telegram",
                          json={"file_id": gen["file_id"]}, headers=_auth(admin))
    assert r.status_code == 200, r.text
    assert r.json()["sent"] is True and sent["chat_id"] == 555
    assert sent["method"] == "sendDocument" and sent["name"].endswith(".pptx")


async def test_send_telegram_rejects_foreign_file(client, monkeypatch):
    admin = await _register_admin(client)
    gen = await _generate(client, admin)
    seller = await _make_user(client, admin, "seller4", "продажи")
    r = await client.post("/api/proposals/send-telegram",
                          json={"file_id": gen["file_id"]}, headers=_auth(seller))
    assert r.status_code == 403                       # чужой файл не отправить


async def test_send_telegram_too_large_returns_link(client, monkeypatch):
    admin = await _register_admin(client)
    gen = await _generate(client, admin)

    import app.routers.proposals as pr
    from app.config import get_settings
    settings = get_settings()
    monkeypatch.setattr(settings, "telegram_bot_token", "test-token", raising=False)
    monkeypatch.setattr(settings, "telegram_allowlist", "555:director", raising=False)
    monkeypatch.setattr(pr, "get_settings", lambda: settings)
    monkeypatch.setattr(pr, "TELEGRAM_FILE_LIMIT", 10)   # искусственно крошечный лимит

    r = await client.post("/api/proposals/send-telegram",
                          json={"file_id": gen["file_id"]}, headers=_auth(admin))
    assert r.status_code == 200
    body = r.json()
    assert body["sent"] is False and body["reason"] == "too_large"
    assert body["download_url"] == f"/api/files/{gen['file_id']}"
