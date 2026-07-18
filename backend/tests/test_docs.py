"""Тесты генерации документов (docgen), редактора docx (doc_editor), извлечения (extract)."""
import io
import re

from services import docgen
from services.extract import detect_kind, extract_text
from utils.doc_editor import apply_doc_edits, read_doc

SAMPLE = "# Договор\n**Жирный** абзац про XCMG.\n\n| A | B |\n|---|---|\n| 1 | 2 |"


def test_create_docx_is_valid():
    from docx import Document
    data = docgen.create_docx(SAMPLE, "test")
    assert len(data) > 1000
    doc = Document(io.BytesIO(data))
    text = "\n".join(p.text for p in doc.paragraphs)
    assert "Договор" in text and "XCMG" in text


def test_create_pdf_is_valid():
    data = docgen.create_pdf("# Заголовок\nКириллический текст.", "test")
    assert data[:5] == b"%PDF-" and len(data) > 800


def test_doc_editor_read_and_edit_roundtrip():
    dx = docgen.create_docx("Первый абзац.\n\nВторой абзац.", "d")
    listing = read_doc(dx)
    assert "P1#" in listing or "параграф" in listing.lower()
    ref = re.search(r"P\d+#\w+", listing).group(0)
    out, diff = apply_doc_edits(dx, [{"op": "rewrite", "ref": ref, "new_text": "Изменённый абзац."}])
    assert len(out) > 1000
    from docx import Document
    text = "\n".join(p.text for p in Document(io.BytesIO(out)).paragraphs)
    assert "Изменённый абзац." in text


def test_create_proposal():
    from docx import Document
    data = {
        "filename": "kp", "title": "Поставка спецтехники", "client": "ООО «Стройка»",
        "markup_percent": 12, "validity_days": 14, "contact": "Иван · +7 900 000",
        "items": [
            {"name": "Экскаватор XCMG XE215C", "qty": 1, "price": 9850000},
            {"name": "Ковш дополнительный", "qty": 2, "price": 150000},
        ],
    }
    out = docgen.create_proposal(data)
    assert len(out) > 1000
    text = "\n".join(p.text for p in Document(io.BytesIO(out)).paragraphs)
    tables_text = " ".join(
        c.text for t in Document(io.BytesIO(out)).tables for row in t.rows for c in row.cells)
    assert "Коммерческое предложение" in text   # заголовок-абзац (не плашка)
    assert "Стройка" in text
    assert "Экскаватор XCMG XE215C" in tables_text
    # экскаватор с наценкой 12%: 9 850 000 * 1.12 = 11 032 000
    assert "11 032 000" in tables_text
    # итог: 11 032 000 + 2*150 000*1.12 = 11 368 000
    assert "11 368 000" in tables_text
    # реквизиты компании присутствуют (issue #26)
    assert "2447007401" in text  # ИНН «Ремтехники»


def test_create_proposal_pdf():
    # issue #28 — КП в PDF
    data = {
        "filename": "kp", "title": "Поставка спецтехники", "client": "ООО «Стройка»",
        "markup_percent": 10, "items": [{"name": "Экскаватор", "qty": 1, "price": 1000000}],
    }
    out = docgen.create_proposal_pdf(data)
    assert out[:5] == b"%PDF-" and len(out) > 1000


def test_create_spec_report():
    # issue #25 — отчёт анализа ТЗ
    from docx import Document
    data = {
        "title": "ТЗ на портал заявок", "summary": "Веб-портал для приёма заявок клиентов.",
        "requirements": ["Авторизация сотрудников", "Личный кабинет клиента"],
        "risks": ["Не указаны сроки"], "contradictions": ["П.3 противоречит П.7"],
        "gaps": ["Не указан объём нагрузки"],
    }
    out = docgen.create_spec_report(data)
    d = Document(io.BytesIO(out))
    text = "\n".join(p.text for p in d.paragraphs)
    tables = " ".join(c.text for t in d.tables for r in t.rows for c in r.cells)
    assert "АНАЛИЗ" in tables and "Требования" in tables
    assert "Авторизация сотрудников" in text and "Не указаны сроки" in text


def test_create_estimate():
    # issue #27 — Excel-смета с настоящими формулами
    from openpyxl import load_workbook
    data = {
        "title": "Смета на ТО", "client": "ООО «Тест»", "markup_percent": 10,
        "items": [
            {"name": "Работа механика", "unit": "ч", "qty": 8, "price": 3600},
            {"name": "Масло моторное", "unit": "л", "qty": 20, "price": 450},
        ],
    }
    out = docgen.create_estimate(data)
    ws = load_workbook(io.BytesIO(out)).active   # формулы как строки
    cells = [str(c.value) for row in ws.iter_rows() for c in row if c.value is not None]
    joined = " ".join(cells)
    assert "Работа механика" in joined and "ИТОГО" in joined
    assert any(v.startswith("=ROUND") for v in cells)   # сумма позиции — формула
    assert any(v.startswith("=SUM(") for v in cells)     # итог — формула


def test_fill_template():
    # issue #26 — подстановка {{ПОЛЕ}} с сохранением структуры
    from docx import Document
    doc = Document()
    doc.add_paragraph("Договор с {{КЛИЕНТ}} на сумму {{ЦЕНА}} рублей.")
    t = doc.add_table(rows=1, cols=2)
    t.rows[0].cells[0].text = "Дата: {{ДАТА}}"
    t.rows[0].cells[1].text = "Поле {{НЕ_ЗАПОЛНЕНО}}"
    buf = io.BytesIO()
    doc.save(buf)

    out, filled, remaining = docgen.fill_template(
        buf.getvalue(), {"КЛИЕНТ": "ООО Ромашка", "ЦЕНА": "1 000 000", "ДАТА": "08.07.2026"})
    d = Document(io.BytesIO(out))
    text = "\n".join(p.text for p in d.paragraphs)
    tables = " ".join(c.text for tb in d.tables for r in tb.rows for c in r.cells)
    assert "ООО Ромашка" in text and "1 000 000" in text and "{{КЛИЕНТ}}" not in text
    assert "08.07.2026" in tables
    assert set(filled) == {"КЛИЕНТ", "ЦЕНА", "ДАТА"}
    assert remaining == ["НЕ_ЗАПОЛНЕНО"]


def test_proposal_has_no_logo():
    # логотип убран из документов — КП не содержит встроенных картинок
    from docx import Document
    out = docgen.create_proposal({"filename": "kp", "items": [{"name": "X", "price": 100}]})
    assert len(Document(io.BytesIO(out)).inline_shapes) == 0


async def test_create_contract_tool(monkeypatch):
    # #43 — договор: реквизиты + правовая оговорка + сохранение полей [УТОЧНИТЬ]
    import app.orchestrator as orch
    captured = {}

    async def fake_save(self, uid, cid, name, data, kind, emit, etype):
        captured["data"] = data

    monkeypatch.setattr(orch.Orchestrator, "_save_file", fake_save)

    async def emit(_e):
        pass
    res = await orch.Orchestrator()._execute_tool(
        "create_contract",
        {"title": "Договор поставки", "filename": "Договор",
         "content": "1. Предмет: поставка экскаватора.\n2. Цена: [УТОЧНИТЬ: сумма договора]."},
        emit, 1, None, None)
    assert "черновик" in res.lower()
    from docx import Document
    text = "\n".join(p.text for p in Document(io.BytesIO(captured["data"])).paragraphs)
    assert "[УТОЧНИТЬ: сумма договора]" in text        # незаполненное сохранено, не выдумано
    assert "ПРАВОВАЯ ОГОВОРКА" in text                 # правовая оговорка добавлена
    assert "2447007401" in text                        # реквизиты «Ремтехники» (ИНН)


def test_create_contract_rbac():
    from agent.registry import role_can_use_tool
    assert role_can_use_tool("продажи", "create_contract") is True
    assert role_can_use_tool("руководство", "create_contract") is True
    assert role_can_use_tool("admin", "create_contract") is True
    assert role_can_use_tool("закупки", "create_contract") is False
    assert role_can_use_tool("user", "create_contract") is False


def test_create_conversation_report():
    # #41 — отчёт анализа переписки/звонка
    from docx import Document
    data = {
        "title": "Звонок с ООО «Ромашка»",
        "summary": "Обсудили поставку экскаватора; клиент готов, ждёт КП.",
        "agreements": ["Поставка LiuGong 975F", "Аванс 50%"],
        "open_questions": ["Точный срок поставки"],
        "next_steps": ["Кирилл: подготовить КП до пятницы"],
        "risks": ["Задержка поставки из Китая"],
    }
    out = docgen.create_conversation_report(data)
    d = Document(io.BytesIO(out))
    text = "\n".join(p.text for p in d.paragraphs)
    tables = " ".join(c.text for t in d.tables for r in t.rows for c in r.cells)
    assert "АНАЛИЗ" in tables and "Договорённости" in tables
    assert "Аванс 50%" in text and "подготовить КП" in text


async def test_analyze_conversation_tool(monkeypatch):
    import app.orchestrator as orch
    captured = {}

    async def fake_save(self, uid, cid, name, data, kind, emit, etype):
        captured["data"] = data

    monkeypatch.setattr(orch.Orchestrator, "_save_file", fake_save)

    async def emit(_e):
        pass
    res = await orch.Orchestrator()._execute_tool(
        "analyze_conversation",
        {"summary": "итог", "agreements": ["a1"], "next_steps": ["s1"], "risks": ["r1"]},
        emit, 1, None, None)
    assert "3 пунктов" in res            # 1 договорённость + 1 шаг + 1 риск
    assert captured["data"][:2] == b"PK"  # валидный .docx (zip)


def test_create_presentation_is_valid():
    # презентация: титул + слайды, тезисы и **жирный** на месте
    from pptx import Presentation
    spec = {
        "title": "Модернизация парка",
        "subtitle": "Для ООО «Стройка» · 2026",
        "slides": [
            {"title": "Проблема", "bullets": ["Износ **65%**", "Простои растут"]},
            {"title": "Решение", "bullets": ["XCMG XE215", "Сервис 24/7"], "notes": "акцент на сроки"},
        ],
    }
    out = docgen.create_presentation(spec)
    assert out[:2] == b"PK" and len(out) > 5000     # валидный .pptx (zip)
    prs = Presentation(io.BytesIO(out))
    assert len(prs.slides._sldIdLst) == 3            # титул + 2 контентных
    all_text = " ".join(
        sh.text_frame.text for s in prs.slides for sh in s.shapes if sh.has_text_frame)
    assert "Модернизация парка" in all_text and "Проблема" in all_text
    assert "65%" in all_text and "XCMG XE215" in all_text
    # заметки докладчика сохранены
    notes = [s.notes_slide.notes_text_frame.text for s in prs.slides if s.has_notes_slide]
    assert any("сроки" in n for n in notes)


async def test_create_presentation_tool(monkeypatch):
    import app.orchestrator as orch
    captured = {}

    async def fake_save(self, uid, cid, name, data, kind, emit, etype):
        captured.update(name=name, data=data, kind=kind)

    monkeypatch.setattr(orch.Orchestrator, "_save_file", fake_save)

    async def emit(_e):
        pass
    res = await orch.Orchestrator()._execute_tool(
        "create_presentation",
        {"title": "Питч", "filename": "Питч",
         "slides": [{"title": "Идея", "bullets": ["раз", "два"]}]},
        emit, 1, None, None)
    assert "2 слайдов" in res                         # титул + 1 контентный
    assert captured["name"] == "Питч.pptx" and captured["kind"] == "pptx"
    assert captured["data"][:2] == b"PK"


def _png_bytes(w=64, h=48, color=(200, 160, 40)):
    from PIL import Image
    buf = io.BytesIO()
    Image.new("RGB", (w, h), color).save(buf, "PNG")
    return buf.getvalue()


def test_presentation_embeds_images_via_bytes():
    # docgen встраивает картинки, переданные байтами (обложка + слайды split/full)
    from pptx import Presentation
    spec = {
        "title": "С картинками", "_cover_image": _png_bytes(160, 90),
        "slides": [
            {"title": "Сплит", "bullets": ["раз"], "_image": _png_bytes(90, 140), "layout": "split"},
            {"title": "Раздел", "bullets": ["акцент"], "_image": _png_bytes(160, 90), "layout": "full"},
            {"title": "Без фото", "bullets": ["текст"]},
        ],
    }
    out = docgen.create_presentation(spec)
    prs = Presentation(io.BytesIO(out))
    pics = sum(1 for s in prs.slides for sh in s.shapes if sh.shape_type == 13)
    assert pics == 3                                   # обложка + 2 слайда с фото


async def test_create_presentation_hybrid_images(monkeypatch):
    # гибрид: реального фото нет → AI-генерация (FLUX замокана) для обложки и слайда
    import app.orchestrator as orch
    from pptx import Presentation
    captured = {}

    async def fake_save(self, uid, cid, name, data, kind, emit, etype):
        captured["data"] = data

    async def fake_flux(prompt, aspect_ratio="1:1"):
        return _png_bytes()

    monkeypatch.setattr(orch.Orchestrator, "_save_file", fake_save)
    monkeypatch.setattr(orch.replicate_svc, "generate_image_flux", fake_flux)
    monkeypatch.setattr(orch.assets, "find_photo", lambda q: None)   # ассетов нет

    events = []

    async def emit(e):
        events.append(e)
    res = await orch.Orchestrator()._execute_tool(
        "create_presentation",
        {"title": "Питч", "cover_image_prompt": "modern excavator, cinematic",
         "slides": [{"title": "Идея", "bullets": ["раз"], "image_prompt": "yellow loader"}]},
        emit, 1, None, None)
    assert "картинок: 2" in res                        # обложка + 1 слайд сгенерированы
    prs = Presentation(io.BytesIO(captured["data"]))
    pics = sum(1 for s in prs.slides for sh in s.shapes if sh.shape_type == 13)
    assert pics == 2
    assert any(e.get("type") == "status" and "изображени" in e.get("text", "") for e in events)


def test_assets_find_photo(tmp_path, monkeypatch):
    # поиск реального фото по ключу (совпадение токенов имени файла)
    from services import assets as assets_mod
    (tmp_path / "xcmg-xe215.jpg").write_bytes(b"JPGDATA")
    (tmp_path / "liugong_856h.png").write_bytes(b"PNGDATA")
    monkeypatch.setattr(assets_mod, "_dir", lambda: str(tmp_path))
    assert assets_mod.find_photo("XCMG XE215") == b"JPGDATA"
    assert assets_mod.find_photo("LiuGong 856H погрузчик") == b"PNGDATA"
    assert assets_mod.find_photo("самосвал КамАЗ") is None      # нет совпадений
    assert set(assets_mod.list_assets()) == {"xcmg-xe215.jpg", "liugong_856h.png"}


def test_create_proposal_pptx_is_valid():
    # #45 — КП-презентация: обложка + слайды блоков + авто-слайд цены; бренд/реквизиты
    from pptx import Presentation
    data = {
        "name": "Экскаватор XCMG XE215C", "brand": "XCMG", "manager": "Иван Петров",
        "client_name": "ООО «Стройка»", "price": "9 850 000 ₽",
        "payment_terms": ["50% аванс", "50% по факту"],
        "blocks": [
            {"type": "title", "title": "Экскаватор XCMG XE215C", "text": "масса 21.5 т"},
            {"type": "split", "rows": [["ДВИГАТЕЛЬ", None], ["Мощность", "118 кВт"]]},
            {"type": "table", "title": "Комплектация", "rows": [["Кондиционер", "есть"]]},
            {"type": "text", "title": "Преимущества", "text": "Надёжный двигатель."},
        ],
    }
    out = docgen.create_proposal_pptx(data)
    assert out[:2] == b"PK" and len(out) > 5000
    prs = Presentation(io.BytesIO(out))
    assert len(prs.slides._sldIdLst) == 5              # 4 блока + слайд цены
    text = " ".join(sh.text_frame.text for s in prs.slides for sh in s.shapes if sh.has_text_frame)
    assert "КОММЕРЧЕСКОЕ ПРЕДЛОЖЕНИЕ" in text and "XCMG XE215C" in text
    assert "118 кВт" in text and "СТОИМОСТЬ" in text and "9 850 000" in text
    assert "2447007401" in text                        # ИНН «Ремтехники» из docx_style


def test_create_proposal_pptx_empty_blocks_has_price():
    # без блоков — всё равно валидный файл со слайдом цены
    from pptx import Presentation
    out = docgen.create_proposal_pptx({"name": "Погрузчик", "blocks": []})
    prs = Presentation(io.BytesIO(out))
    assert len(prs.slides._sldIdLst) == 1              # только авто-слайд цены


async def test_create_proposal_pptx_tool(monkeypatch):
    # проводка инструмента + резолв image_asset из папки-ассетов
    import app.orchestrator as orch
    from pptx import Presentation
    captured = {}

    async def fake_save(self, uid, cid, name, data, kind, emit, etype):
        captured.update(name=name, data=data, kind=kind)

    monkeypatch.setattr(orch.Orchestrator, "_save_file", fake_save)
    monkeypatch.setattr(orch.assets, "find_photo", lambda k: _png_bytes() if k else None)

    async def emit(_e):
        pass
    res = await orch.Orchestrator()._execute_tool(
        "create_proposal_pptx",
        {"name": "Экскаватор", "filename": "КП_XE215",
         "blocks": [{"type": "title", "title": "Экскаватор"},
                    {"type": "split", "image_asset": "XCMG XE215",
                     "rows": [["Мощность", "118 кВт"]]}]},
        emit, 1, None, None)
    assert "3 слайдов" in res                           # 2 блока + слайд цены
    assert captured["name"] == "КП_XE215.pptx" and captured["kind"] == "pptx"
    prs = Presentation(io.BytesIO(captured["data"]))
    pics = sum(1 for s in prs.slides for sh in s.shapes if sh.shape_type == 13)
    assert pics == 1                                    # фото подставлено из ассетов


def test_create_proposal_pptx_rbac():
    from agent.registry import role_can_use_tool
    assert role_can_use_tool("продажи", "create_proposal_pptx") is True
    assert role_can_use_tool("руководство", "create_proposal_pptx") is True
    assert role_can_use_tool("admin", "create_proposal_pptx") is True
    assert role_can_use_tool("закупки", "create_proposal_pptx") is False


def test_detect_kind():
    assert detect_kind("a.docx") == "docx"
    assert detect_kind("b.PDF") == "pdf"
    assert detect_kind("c.png") == "image"
    assert detect_kind("d.xlsx") == "xlsx"
    assert detect_kind("call.mp3") == "audio"
    assert detect_kind("rec.ogg") == "audio"
    assert detect_kind("e.unknown") == "other"


def test_extract_text_from_docx():
    dx = docgen.create_docx("Прайс на запчасти XCMG.", "p")
    text = extract_text(dx, "p.docx")
    assert "запчасти" in text
